# -*- coding: utf-8 -*-
"""Bharat Direct pitch deck (.pptx) — 16:9."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x11, 0x18, 0x3A)
NAVY2 = RGBColor(0x1B, 0x25, 0x55)
BLUE = RGBColor(0x3A, 0x5B, 0xA9)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
AMBER = RGBColor(0xF4, 0xA8, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x60, 0x70)
LGREY = RGBColor(0xEE, 0xF1, 0xF7)
DARKTX = RGBColor(0x1A, 0x1F, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(blank)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def tb(s, l,t,w,h, text, size=18, color=DARKTX, bold=False, italic=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri', line=1.05, space=6):
    box = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    lines = text.split('\n')
    for i,ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line; p.space_after = Pt(space)
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.name = font; f.color.rgb = color
    return box

def bullets(s, l,t,w,h, items, size=16, color=DARKTX, gap=8, lead_color=AMBER):
    box = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,(lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing = 1.05; p.space_after = Pt(gap)
        r0 = p.add_run(); r0.text = '› '; r0.font.size=Pt(size); r0.font.bold=True; r0.font.color.rgb=lead_color; r0.font.name='Calibri'
        if lead:
            r1 = p.add_run(); r1.text = lead; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=color; r1.font.name='Calibri'
        r2 = p.add_run(); r2.text = rest; r2.font.size=Pt(size); r2.font.color.rgb=color; r2.font.name='Calibri'
    return box

def circle_num(s, l, t, d, n, fill=AMBER, txt=NAVY):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l),Inches(t),Inches(d),Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background(); c.shadow.inherit=False
    tf = c.text_frame; tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=txt; r.font.name='Calibri'
    return c

def card(s, l,t,w,h, fill=LGREY):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.fill.background(); box.shadow.inherit=False
    return box

def kicker(s, text, color=AMBER, l=0.85, t=0.55):
    tb(s, l,t,11,0.4, text.upper(), size=13, color=color, bold=True)

def title(s, text, l=0.85, t=0.95, w=11.6, size=33, color=NAVY):
    tb(s, l,t,w,1.2, text, size=size, color=color, bold=True, line=1.0)

# ---- 1 TITLE ----
s = slide(NAVY)
tb(s, 0.85, 2.2, 11.6, 1.4, 'BHARAT DIRECT', size=60, color=WHITE, bold=True)
tb(s, 0.9, 3.55, 11.6, 0.7, 'Finishing what UPI started', size=24, color=AMBER, italic=True)
tb(s, 0.9, 4.5, 11.0, 1.2, 'A bank-owned, aggregator-free payments rail — and a Fraud\nIntelligence Consortium to close the human gap UPI left open.', size=18, color=ICE, line=1.2)
tb(s, 0.9, 6.6, 11.0, 0.5, 'Discussion paper for regulators, banks and the payments industry  ·  Working draft', size=12, color=ICE)

# ---- 2 TWO PROBLEMS ----
s = slide()
kicker(s, 'The problem')
title(s, 'India built 80%. Two things remain.')
card(s, 0.85, 2.3, 5.7, 3.9, LGREY)
circle_num(s, 1.2, 2.65, 0.8, '₹', fill=AMBER)
tb(s, 1.2, 3.65, 5.0, 0.6, 'The toll-booth', size=22, color=NAVY, bold=True)
tb(s, 1.2, 4.3, 5.05, 1.7, 'A private payment-aggregator layer still collects rent between merchant and bank — even though the money can move bank-to-bank.', size=15, color=DARKTX, line=1.15)
card(s, 6.75, 2.3, 5.7, 3.9, NAVY)
circle_num(s, 7.1, 2.65, 0.8, '!', fill=AMBER)
tb(s, 7.1, 3.65, 5.0, 0.6, 'The fraud gap', size=22, color=WHITE, bold=True)
tb(s, 7.1, 4.3, 5.05, 1.7, 'Social-engineering (APP) fraud — a victim manipulated into authorising a real payment. The one gap a faster rail does nothing to close.', size=15, color=ICE, line=1.15)
tb(s, 0.85, 6.5, 11.6, 0.5, 'Pillar 1 removes the rent.   Pillar 2 removes the fear.', size=17, color=BLUE, bold=True)

# ---- 3 THESIS ----
s = slide(NAVY)
kicker(s, 'The thesis', color=AMBER)
tb(s, 0.85, 1.2, 11.6, 2.0, 'This is not a new paradigm.\nIt completes a proven one.', size=34, color=WHITE, bold=True, line=1.05)
items = [
 ('NPCI ', 'is already a bank-owned consortium.'),
 ('CERSAI / CKYC ', 'is already a central identity utility.'),
 ('UPI ', 'already moves money bank-to-bank with no aggregator for person-to-merchant payments.'),
]
bullets(s, 0.9, 3.7, 11.0, 2.0, items, size=19, color=ICE, gap=14, lead_color=AMBER)
tb(s, 0.9, 6.3, 11.4, 0.8, 'Bharat Direct extends that co-owned model up into merchant services — and adds the one shared utility India hasn’t built yet: fraud intelligence.', size=16, color=WHITE, italic=True, line=1.15)

# ---- 4 PSYCHOLOGY ----
s = slide()
kicker(s, 'Why people who "know better" still pay')
title(s, 'A scam hijacks the brain, not the knowledge')
tb(s, 0.85, 2.15, 11.5, 0.8, 'Victims almost always know scams exist. Education fails because knowledge stored yesterday is unavailable in the moment of attack.', size=16, color=GREY, line=1.15)
cards = [
 ('Hot state','Panic, fear or urgency narrows judgment. The calm "I’d never fall for it" person and the panicked victim are different decision-makers.'),
 ('Pressure levers','Authority, urgency, reward, trust, and a small first step that escalates.'),
 ('Isolation + tempo','"Don’t hang up. Don’t tell anyone." A closed loop with one voice — no second opinion can enter.'),
]
x = 0.85
for i,(h,b) in enumerate(cards):
    card(s, x, 3.2, 3.75, 2.9, LGREY)
    circle_num(s, x+0.3, 3.5, 0.55, i+1)
    tb(s, x+0.3, 4.2, 3.2, 0.5, h, size=17, color=NAVY, bold=True)
    tb(s, x+0.3, 4.75, 3.25, 1.2, b, size=12.5, color=DARKTX, line=1.1)
    x += 4.0

# ---- 5 CHOKEPOINT ----
s = slide(NAVY)
kicker(s, 'The design conclusion', color=AMBER)
tb(s, 0.85, 1.4, 11.6, 3.0, 'Every scam converges on ONE moment:\nthe victim authorising the payment.', size=34, color=WHITE, bold=True, line=1.1)
tb(s, 0.9, 3.9, 11.2, 1.6, 'You can’t stop people being called. You can only intervene at the moment of authorisation — break the closed loop, cool the hot state, and bring deliberate thinking back online for two seconds.', size=19, color=ICE, line=1.25)
tb(s, 0.9, 6.2, 11.2, 0.6, 'Everything before it is manipulation. Everything after it is irreversible.', size=16, color=AMBER, italic=True, bold=True)

# ---- 6 WHY ADS FAIL ----
s = slide()
kicker(s, 'Why scattered ads don’t work')
title(s, 'Habituation kills the generic warning')
tb(s, 0.85, 2.2, 11.4, 1.0, 'A "Beware of fraud" banner on every transaction becomes wallpaper within a week. The design law:', size=17, color=GREY, line=1.15)
card(s, 0.85, 3.2, 11.6, 1.5, NAVY)
tb(s, 1.2, 3.45, 10.9, 1.0, '"Warnings must be rare, specific, and earned by real risk — or they are worse than nothing."', size=22, color=WHITE, bold=True, italic=True, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 0.85, 5.1, 11.5, 1.6, 'That is exactly why a Fraud Intelligence Consortium is not decoration: only shared, precise intelligence lets a warning fire ONLY when it matters — which is the only way it keeps its power.', size=18, color=DARKTX, line=1.25)

# ---- 7 AWARE LAYER ----
s = slide()
kicker(s, 'Pillar 2a — the screen')
title(s, 'The Aware Layer: spot-awareness at the chokepoint')
rows = [
 ('Dynamic friction','Frictionless normally; staged interruption only on high risk.'),
 ('Cooling interstitial','Unskippable countdown on a high-risk new payee — breaks urgency.'),
 ('Trance-breaker','"Is someone on the phone right now telling you to pay?"'),
 ('Named-threat warning','"This matches a digital-arrest scam." Specific, not generic.'),
 ('Payee reputation','"This account has 9 fraud reports in 24h" — the missing second opinion.'),
 ('Active-recall confirm','Type the purpose; re-read the payee. Defeats autopilot.'),
]
x,y = 0.85, 2.2
for i,(h,b) in enumerate(rows):
    col = i % 2; row = i // 2
    cx = 0.85 + col*6.0; cy = 2.2 + row*1.55
    card(s, cx, cy, 5.7, 1.35, LGREY)
    circle_num(s, cx+0.25, cy+0.27, 0.5, i+1)
    tb(s, cx+0.95, cy+0.18, 4.6, 0.45, h, size=15, color=NAVY, bold=True)
    tb(s, cx+0.95, cy+0.62, 4.65, 0.7, b, size=11.5, color=DARKTX, line=1.05)
tb(s, 0.85, 7.0, 11.6, 0.4, 'The screen becomes the second voice in the room.', size=14, color=BLUE, italic=True, bold=True)

# ---- 8 FIC ----
s = slide(NAVY)
kicker(s, 'Pillar 2b — the brain', color=AMBER)
title(s, 'The Fraud Intelligence Consortium', color=WHITE)
tb(s, 0.85, 2.05, 11.5, 0.9, 'Fraud is a network crime fought by isolated institutions. A mule drains victims across ten banks; each sees only its sliver. The criminals live in the gaps.', size=15, color=ICE, line=1.15)
items = [
 ('Mule / payee graph — ','a report at one bank reaches another’s confirm-screen in seconds.'),
 ('Scam-typology library — ','powers the named-threat warnings.'),
 ('Cross-bank layering detection — ','follows money no single bank can see.'),
 ('Federated ML — ','signal shared, raw customer data never. Privacy- and competition-safe.'),
 ('Golden-hour freeze — ','instant cross-bank freeze and recovery.'),
]
bullets(s, 0.9, 3.2, 11.2, 2.6, items, size=15.5, color=ICE, gap=8, lead_color=AMBER)
tb(s, 0.9, 6.45, 11.4, 0.7, 'Incentive backbone: make the receiving bank share APP liability. Then contributing to the FIC becomes self-interest, not charity.', size=14, color=WHITE, italic=True, line=1.1)

# ---- 9 PILLAR 1 + ARCH ----
s = slide()
kicker(s, 'Pillar 1 — the rail')
title(s, 'Direct Settlement Consortium')
tb(s, 0.85, 2.1, 11.4, 1.0, 'Merchants connect directly to banks. Funds settle bank-to-bank over a thin switch that passes messages but never holds money — so the entire PA escrow regime disappears.', size=16, color=GREY, line=1.2)
# simple architecture row
labels = [('CKYC / Identity','one-time, central'),('Thin switch + netting','no fund custody'),('Banks own acquiring','the "pieces"')]
x = 0.95
for i,(h,b) in enumerate(labels):
    card(s, x, 3.4, 3.6, 1.7, NAVY if i==1 else LGREY)
    tb(s, x+0.25, 3.7, 3.1, 0.6, h, size=16, color=(WHITE if i==1 else NAVY), bold=True, line=1.0)
    tb(s, x+0.25, 4.45, 3.1, 0.5, b, size=12, color=(ICE if i==1 else GREY))
    if i<2:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+3.65), Inches(4.0), Inches(0.45), Inches(0.45))
        ar.fill.solid(); ar.fill.fore_color.rgb = AMBER; ar.line.fill.background(); ar.shadow.inherit=False
    x += 4.05
tb(s, 0.85, 5.6, 11.5, 1.2, 'Only identity stays central, amortised to near-zero per transaction. The percentage merchant fee becomes a small flat per-message cost.', size=16, color=DARKTX, line=1.2)

# ---- 10 WHO GAINS ----
s = slide()
kicker(s, 'Economics')
title(s, 'Who gains')
data = [
 ('Small / mid merchants','Flat near-zero cost vs a percentage fee'),
 ('Banks','Value-added revenue PAs take today + a flat regulated rail fee + lower fraud losses'),
 ('Government','Lower subsidy bill · citizen protection · sovereignty · forex retained'),
 ('Consumers','Lower prices and far lower scam risk'),
 ('Fintechs / PAs','Repositioned as builders of the Aware Layer & FIC — not killed'),
]
y = 2.2
for i,(a,b) in enumerate(data):
    fill = LGREY if i%2==0 else WHITE
    card(s, 0.85, y, 11.6, 0.82, fill)
    tb(s, 1.15, y+0.16, 3.3, 0.5, a, size=15, color=NAVY, bold=True)
    tb(s, 4.6, y+0.16, 7.6, 0.55, b, size=14, color=DARKTX, line=1.0)
    y += 0.9
tb(s, 0.85, 6.85, 11.6, 0.5, 'Keystone: the rail charges a small flat regulated fee — not zero. That is what makes it self-funding.', size=14, color=AMBER, bold=True)

# ---- 11 NUMBERS ----
s = slide(NAVY)
kicker(s, 'The numbers', color=AMBER)
title(s, 'The scale is already here', color=WHITE)
stats = [
 ('21.7 bn','UPI transactions in a single month (Jan 2026)'),
 ('₹28.3 L cr','value moved that month'),
 ('₹0','MDR on UPI / RuPay debit since 2020'),
 ('~₹2,000 cr','annual govt subsidy to keep it free'),
]
x = 0.85
for i,(n,l) in enumerate(stats):
    card(s, x, 2.6, 2.85, 2.6, NAVY2)
    tb(s, x+0.2, 3.0, 2.5, 0.9, n, size=33, color=AMBER, bold=True)
    tb(s, x+0.2, 4.0, 2.5, 1.1, l, size=13, color=ICE, line=1.1)
    x += 3.0
tb(s, 0.85, 5.7, 11.6, 0.9, 'India spends real money keeping payments free. A co-owned, lower-cost rail shrinks that bill — and a fraud consortium protects the citizens using it.', size=16, color=WHITE, line=1.2)

# ---- 12 GOVERNMENT ----
s = slide()
kicker(s, 'The pitch to government')
title(s, 'Three sovereign wins')
gov = [
 ('Fiscal','The State already spends ₹1,500–2,000 cr/yr subsidising UPI. A lower-cost rail shrinks that bill.'),
 ('Citizen protection','Driving scam losses toward zero is rare, tangible political capital — protection people feel.'),
 ('Strategic autonomy','Scheme fees leave the country; mule networks exploit gaps. Both fixed by infrastructure India owns.'),
]
x = 0.85
for i,(h,b) in enumerate(gov):
    card(s, x, 2.4, 3.75, 3.3, LGREY)
    circle_num(s, x+0.3, 2.7, 0.6, i+1)
    tb(s, x+0.3, 3.5, 3.2, 0.6, h, size=18, color=NAVY, bold=True, line=1.0)
    tb(s, x+0.3, 4.2, 3.25, 1.4, b, size=13, color=DARKTX, line=1.15)
    x += 4.0
tb(s, 0.85, 6.1, 11.6, 0.9, 'India built the hardest 80%. Bharat Direct finishes it — removing the last private toll-booth and closing the one gap UPI left open: the human at the moment of payment.', size=16, color=BLUE, italic=True, bold=True, line=1.2)

# ---- 13 RISKS ----
s = slide()
kicker(s, 'Honest risks')
title(s, 'Named before they’re asked')
risks = [
 ('Concentration','A shared switch is a systemic target → federated, multi-operator design, not a monolith.'),
 ('Antitrust','A consortium mustn’t be a cartel → neutral not-for-profit governance, open membership (NPCI precedent).'),
 ('Privacy','No pooling of raw data → federated learning by architecture.'),
 ('Transition cost','PA/fintech jobs and tax → repositioning narrative is essential.'),
 ('Residual APP fraud','Won’t hit literal zero → target world-best, backed by reimbursement so the victim is made whole.'),
]
bullets(s, 0.85, 2.3, 11.6, 4.0, [(h+' — ', b) for h,b in risks], size=17, color=DARKTX, gap=14, lead_color=AMBER)

# ---- 14 CALL TO ACTION ----
s = slide(NAVY)
kicker(s, 'Call to action', color=AMBER)
tb(s, 0.85, 1.15, 11.6, 1.4, 'A pre-competitive working group —\nthe way UPI itself began.', size=32, color=WHITE, bold=True, line=1.05)
who = [
 ('Regulators — ','RBI (DPSS), MeitY, I4C / cybercrime (1930), DFS'),
 ('Rails — ','NPCI, CERSAI'),
 ('Banks — ','4–6 anchor public & private banks as founding owners'),
 ('Industry — ','Payments Council of India + fintechs as builders'),
 ('Academia — ','behavioural-science & security researchers to measure it'),
]
bullets(s, 0.9, 3.3, 11.2, 2.6, who, size=16, color=ICE, gap=8, lead_color=AMBER)
tb(s, 0.9, 6.3, 11.4, 0.9, 'First deliverable: pilot the Aware Layer + a minimal mule-graph on existing UPI rails — prove the fraud number before touching settlement. Lead with the fraud win.', size=15, color=WHITE, italic=True, line=1.15)

prs.save(r'C:\Users\dhire\OneDrive\Documents\00 Dhirendra\Innovations\Payment Innovation\Bharat Direct\Bharat-Direct-Deck.pptx')
print('Saved deck with', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
